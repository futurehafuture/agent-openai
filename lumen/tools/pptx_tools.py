"""PowerPoint generation tool built on python-pptx.

Produces a clean, themed 16:9 deck from a structured outline. Slides may carry
bullet points, a paragraph body, an embedded image (e.g. a chart produced by the
data tools), and speaker notes — so "analyse data → make charts → build a deck"
composes naturally.
"""

from __future__ import annotations

import json
from dataclasses import dataclass

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from ..logging_setup import get_logger
from ..workspace import workspace
from .registry import register_tool

logger = get_logger(__name__)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


@dataclass(frozen=True)
class Theme:
    bg: RGBColor
    panel: RGBColor
    title: RGBColor
    body: RGBColor
    muted: RGBColor
    accent: RGBColor
    font_title: str
    font_body: str


_THEMES = {
    "light": Theme(
        bg=RGBColor(0xFA, 0xFB, 0xFC),
        panel=RGBColor(0xFF, 0xFF, 0xFF),
        title=RGBColor(0x14, 0x1B, 0x2E),
        body=RGBColor(0x3A, 0x41, 0x52),
        muted=RGBColor(0x8A, 0x90, 0x9E),
        accent=RGBColor(0x3B, 0x5B, 0xDB),
        font_title="Helvetica Neue",
        font_body="Helvetica Neue",
    ),
    "dark": Theme(
        bg=RGBColor(0x12, 0x14, 0x1A),
        panel=RGBColor(0x1B, 0x1E, 0x27),
        title=RGBColor(0xF4, 0xF6, 0xFB),
        body=RGBColor(0xC5, 0xCA, 0xD6),
        muted=RGBColor(0x7C, 0x82, 0x92),
        accent=RGBColor(0x6E, 0x8B, 0xFF),
        font_title="Helvetica Neue",
        font_body="Helvetica Neue",
    ),
    "warm": Theme(
        bg=RGBColor(0xFB, 0xF7, 0xF1),
        panel=RGBColor(0xFF, 0xFD, 0xFA),
        title=RGBColor(0x37, 0x2A, 0x1E),
        body=RGBColor(0x5A, 0x4B, 0x3C),
        muted=RGBColor(0x9C, 0x8B, 0x77),
        accent=RGBColor(0xC9, 0x6A, 0x2E),
        font_title="Georgia",
        font_body="Helvetica Neue",
    ),
}


def _fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_bg(slide, theme: Theme) -> None:
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _fill(rect, theme.bg)
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(2, rect._element)


def _text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int,
    color: RGBColor,
    font: str,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color
    return box


def _bullets(slide, left, top, width, height, items: list[str], theme: Theme) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        dot = p.add_run()
        dot.text = "•  "
        dot.font.color.rgb = theme.accent
        dot.font.size = Pt(18)
        dot.font.bold = True
        run = p.add_run()
        run.text = str(item)
        run.font.size = Pt(18)
        run.font.name = theme.font_body
        run.font.color.rgb = theme.body


def _title_slide(prs, theme: Theme, title: str, subtitle: str | None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, theme)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.7), Inches(0.9), Inches(0.12))
    _fill(bar, theme.accent)
    _text(slide, Inches(0.85), Inches(2.95), Inches(11.6), Inches(2.0), title,
          size=44, color=theme.title, font=theme.font_title, bold=True)
    if subtitle:
        _text(slide, Inches(0.9), Inches(4.3), Inches(11.5), Inches(1.0), subtitle,
              size=20, color=theme.muted, font=theme.font_body)


def _content_slide(prs, theme: Theme, spec: dict, index: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, theme)
    title = str(spec.get("title", f"Slide {index}"))
    _text(slide, Inches(0.85), Inches(0.6), Inches(11.6), Inches(1.0), title,
          size=30, color=theme.title, font=theme.font_title, bold=True)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(1.5), Inches(0.6), Inches(0.06))
    _fill(accent, theme.accent)

    image = spec.get("image")
    has_image = bool(image)
    content_w = Inches(6.0) if has_image else Inches(11.6)

    bullets = spec.get("bullets")
    body = spec.get("body")
    if bullets:
        _bullets(slide, Inches(0.9), Inches(1.9), content_w, Inches(4.8),
                 [str(b) for b in bullets], theme)
    elif body:
        _text(slide, Inches(0.9), Inches(1.9), content_w, Inches(4.8), str(body),
              size=18, color=theme.body, font=theme.font_body)

    if has_image:
        img_path = workspace.resolve_read(str(image))
        if img_path.exists():
            slide.shapes.add_picture(str(img_path), Inches(7.1), Inches(1.9), width=Inches(5.4))
        else:
            _text(slide, Inches(7.1), Inches(3.0), Inches(5.4), Inches(0.6),
                  f"[missing image: {image}]", size=12, color=theme.muted, font=theme.font_body)

    notes = spec.get("notes")
    if notes:
        slide.notes_slide.notes_text_frame.text = str(notes)

    _text(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4), str(index),
          size=11, color=theme.muted, font=theme.font_body, align=PP_ALIGN.RIGHT)


def create_presentation(
    title: str,
    slides_json: str,
    subtitle: str | None = None,
    theme: str = "light",
    output_name: str | None = None,
) -> str:
    """Create a polished .pptx presentation and save it to the workspace.

    Args:
        title: The deck title (shown on the cover slide).
        slides_json: A JSON array of slide objects. Each object supports:
            ``title`` (str), ``bullets`` (list of str) OR ``body`` (str),
            ``image`` (str: a path to an image/chart in the workspace, optional),
            and ``notes`` (str: speaker notes, optional).
            Example: ``[{"title": "Results", "bullets": ["Revenue up 12%",
            "Churn down"], "image": "charts/bar-region.png"}]``
        subtitle: Optional cover-slide subtitle.
        theme: Visual theme: 'light', 'dark', or 'warm'.
        output_name: Optional output filename (defaults to a slug of the title).
    """
    theme_obj = _THEMES.get(theme.lower(), _THEMES["light"])
    try:
        slides = json.loads(slides_json)
    except json.JSONDecodeError as exc:
        raise ValueError(f"slides_json is not valid JSON: {exc}") from exc
    if not isinstance(slides, list) or not slides:
        raise ValueError("slides_json must be a non-empty JSON array of slide objects.")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    _title_slide(prs, theme_obj, title, subtitle)
    for i, spec in enumerate(slides, start=1):
        if not isinstance(spec, dict):
            raise ValueError(f"Slide {i} must be an object, got {type(spec).__name__}.")
        _content_slide(prs, theme_obj, spec, i)

    slug = (output_name or title or "presentation").strip()
    if not slug.lower().endswith(".pptx"):
        slug = slug.replace(" ", "-").lower() + ".pptx"
    target = workspace.unique_output(slug)
    prs.save(str(target))
    return (
        f"📑 Created a {len(slides) + 1}-slide deck (theme: {theme}) → "
        f"**{workspace.display(target)}**."
    )


register_tool(create_presentation, category="skills")
